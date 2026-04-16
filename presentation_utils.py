# presentation_utils.py
import requests
from lxml import etree as ET
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

# Modern Academic Light Aesthetics
PALETTE = {
    "bg": "F8F9FA",          # Light Gray/White
    "text_light": "212529",  # Dark Text 
    "accent": "00509E",      # Corporate Blue
    "title_text": "003366"   # Darker Navy
}

def download_image(url, filename, cache_dir="images"):
    if not url.startswith('http'):
        return url
        
    if not os.path.exists(cache_dir):
        os.makedirs(cache_dir)
        
    filepath = os.path.join(cache_dir, filename)
    if os.path.exists(filepath):
        return filepath # Return cached picture
        
    print(f"Downloading {filename}...")
    try:
        response = requests.get(url, timeout=10)
        response.raise_for_status()
        with open(filepath, 'wb') as f:
            f.write(response.content)
        return filepath
    except Exception as e:
        print(f"Failed to download {filename}: {e}")
        return None

def set_slide_background_color(slide, hex_color_str):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(hex_color_str)

def apply_fade_transition(slide):
    # XML injection for Fade Transition across slides
    sld = slide.element
    if sld.xpath('./p:transition'):
        return
    transition = ET.Element('{http://schemas.openxmlformats.org/presentationml/2006/main}transition', spd="med")
    fade = ET.SubElement(transition, '{http://schemas.openxmlformats.org/presentationml/2006/main}fade')
    sld.insert(0, transition)

def format_title_text(title_shape, text):
    title_shape.text = text
    for paragraph in title_shape.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.LEFT
        for run in paragraph.runs:
            run.font.color.rgb = RGBColor.from_string(PALETTE["accent"])
            run.font.name = 'Helvetica Neue'
            run.font.bold = True

def add_bullets_to_frame(text_frame, bullet_texts):
    for i, b_text in enumerate(bullet_texts):
        if i == 0:
            p = text_frame.paragraphs[0]
            p.text = b_text
        else:
            p = text_frame.add_paragraph()
            p.text = b_text
            
        p.alignment = PP_ALIGN.LEFT
        for run in p.runs:
            run.font.color.rgb = RGBColor.from_string(PALETTE["text_light"])
            run.font.name = 'Helvetica'
            # Differentiate main bullets vs sub points based on indentations
            if "  -" in b_text or "  1" in b_text or "  2" in b_text:
                p.level = 1
                run.font.size = Pt(16)
            else:
                p.level = 0
                run.font.size = Pt(20)
