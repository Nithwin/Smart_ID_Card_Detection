#!/usr/bin/env python3
"""
Premium IEEE Conference Presentation Generator
Smart ID Card Detection and Authentication Using Deep Learning and Computer Vision

Modules used:
  - python-pptx: Core PPTX generation
  - lxml: XML manipulation for transitions
  - Pillow (PIL): Image processing
  - os/copy: File system operations
"""

import os
import copy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from lxml import etree

# ═══════════════════════════════════════════════════════════════
#  DESIGN SYSTEM — Premium Modern Dark Theme
# ═══════════════════════════════════════════════════════════════

# Slide dimensions (standard widescreen 16:9)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Color Palette — Modern Dark Gradient
BG_DARK       = "0F1624"   # Deep space navy
BG_CARD       = "1A2332"   # Slightly lighter card background
ACCENT_CYAN   = "00D4FF"   # Electric cyan accent
ACCENT_BLUE   = "4A90D9"   # Soft blue
ACCENT_PURPLE = "7C3AED"   # Vibrant purple
TEXT_WHITE     = "FFFFFF"   # Pure white
TEXT_MUTED     = "94A3B8"   # Muted silver
TEXT_HEADING   = "E2E8F0"   # Off-white for headings
GRADIENT_START = "1E3A5F"   # Gradient bar start
GRADIENT_END   = "0F1624"   # Gradient bar end
ACCENT_GREEN   = "10B981"   # Success green
ACCENT_ORANGE  = "F59E0B"   # Warning amber

# Typography
FONT_TITLE    = "Segoe UI"
FONT_BODY     = "Segoe UI"
FONT_MONO     = "Consolas"

# Layout Constants (pixel-perfect alignment grid)
MARGIN_LEFT   = Inches(0.8)
MARGIN_RIGHT  = Inches(0.8)
MARGIN_TOP    = Inches(0.6)
CONTENT_TOP   = Inches(1.6)
CONTENT_WIDTH = Inches(11.733)  # Full width minus margins
HALF_WIDTH    = Inches(5.6)
IMAGE_RIGHT_X = Inches(7.8)
IMAGE_WIDTH   = Inches(4.8)

# Image paths
IMG_DIR = "/home/shadow/.gemini/antigravity/brain/ad925e9c-6010-436d-9765-055f24e05bc7"
IMAGES = {
    "yolo":       os.path.join(IMG_DIR, "yolo_id_card_1775276953932.png"),
    "ocr":        os.path.join(IMG_DIR, "ocr_extraction_1775276975776.png"),
    "face":       os.path.join(IMG_DIR, "face_matching_1775276992948.png"),
    "pipeline":   os.path.join(IMG_DIR, "pipeline_architecture_1775277428249.png"),
    "spoof":      os.path.join(IMG_DIR, "anti_spoofing_1775277442496.png"),
    "results":    os.path.join(IMG_DIR, "results_chart_1775277463977.png"),
    "deployment": os.path.join(IMG_DIR, "field_deployment_1775277482929.png"),
}

# ═══════════════════════════════════════════════════════════════
#  HELPER FUNCTIONS
# ═══════════════════════════════════════════════════════════════

def set_slide_bg(slide, hex_color):
    """Set solid background color for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(hex_color)

def add_fade_transition(slide):
    """Inject smooth fade transition via XML."""
    from pptx.oxml.ns import qn
    sld = slide.element
    # Check if transition already exists
    existing = sld.findall(qn('p:transition'))
    if existing:
        return
    transition = etree.SubElement(sld, qn('p:transition'))
    transition.set('spd', 'med')
    transition.set('advClick', '1')
    etree.SubElement(transition, qn('p:fade'))

def add_textbox(slide, left, top, width, height):
    """Add a textbox shape and return text_frame."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    return tf

def add_styled_text(tf, text, font_size=18, color=TEXT_WHITE, bold=False,
                    font_name=FONT_BODY, alignment=PP_ALIGN.LEFT, spacing_after=Pt(6)):
    """Add a paragraph with precise styling to a text_frame."""
    p = tf.add_paragraph()
    p.alignment = alignment
    p.space_after = spacing_after
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = RGBColor.from_string(color)
    run.font.bold = bold
    run.font.name = font_name
    return p

def set_first_paragraph(tf, text, font_size=18, color=TEXT_WHITE, bold=False,
                        font_name=FONT_BODY, alignment=PP_ALIGN.LEFT):
    """Style the default first paragraph of a text_frame."""
    p = tf.paragraphs[0]
    p.alignment = alignment
    p.space_after = Pt(6)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = RGBColor.from_string(color)
    run.font.bold = bold
    run.font.name = font_name
    return p

def add_accent_bar(slide, left, top, width, height, color=ACCENT_CYAN):
    """Add a thin colored accent bar/rectangle."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(color)
    shape.line.fill.background()
    return shape

def add_slide_number(slide, number, total=15):
    """Add slide number indicator at bottom right."""
    tf = add_textbox(slide, Inches(11.5), Inches(6.9), Inches(1.5), Inches(0.4))
    set_first_paragraph(tf, f"{number} / {total}", font_size=10, color=TEXT_MUTED,
                        alignment=PP_ALIGN.RIGHT)

def add_section_title(slide, title_text, subtitle_text=None):
    """Add the main section title with accent bar."""
    # Accent bar
    add_accent_bar(slide, MARGIN_LEFT, Inches(0.5), Inches(0.15), Inches(0.55))

    # Title
    tf = add_textbox(slide, Inches(1.15), Inches(0.4), Inches(10), Inches(0.7))
    set_first_paragraph(tf, title_text, font_size=28, color=TEXT_WHITE,
                        bold=True, font_name=FONT_TITLE)

    # Subtitle / thin divider line
    add_accent_bar(slide, MARGIN_LEFT, Inches(1.15), Inches(3), Inches(0.02), ACCENT_CYAN)

    if subtitle_text:
        tf2 = add_textbox(slide, MARGIN_LEFT, Inches(1.2), Inches(10), Inches(0.4))
        set_first_paragraph(tf2, subtitle_text, font_size=14, color=TEXT_MUTED,
                            font_name=FONT_BODY)

def add_image_safe(slide, img_key, left, top, width=None, height=None):
    """Safely add an image if it exists."""
    path = IMAGES.get(img_key, "")
    if path and os.path.exists(path):
        kwargs = {}
        if width:
            kwargs['width'] = width
        if height:
            kwargs['height'] = height
        slide.shapes.add_picture(path, left, top, **kwargs)
        return True
    return False

def add_bullet_block(slide, bullets, left, top, width, start_size=17, color=TEXT_WHITE):
    """Add a block of bullet points with consistent styling."""
    tf = add_textbox(slide, left, top, width, Inches(5))
    for i, bullet in enumerate(bullets):
        if i == 0:
            set_first_paragraph(tf, bullet, font_size=start_size, color=color)
        else:
            add_styled_text(tf, bullet, font_size=start_size, color=color,
                            spacing_after=Pt(10))
    return tf

# ═══════════════════════════════════════════════════════════════
#  SLIDE BUILDERS
# ═══════════════════════════════════════════════════════════════

def build_title_slide(prs):
    """Slide 1: Title Slide — Dramatic centered layout."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_bg(slide, BG_DARK)
    add_fade_transition(slide)

    # Top accent line
    add_accent_bar(slide, Inches(4), Inches(1.8), Inches(5.333), Inches(0.04), ACCENT_CYAN)

    # Main title
    tf = add_textbox(slide, Inches(1.5), Inches(2.1), Inches(10.333), Inches(1.2))
    set_first_paragraph(tf, "Smart ID Card Detection", font_size=40,
                        color=TEXT_WHITE, bold=True, font_name=FONT_TITLE,
                        alignment=PP_ALIGN.CENTER)
    add_styled_text(tf, "and Authentication", font_size=40, color=ACCENT_CYAN,
                    bold=True, font_name=FONT_TITLE, alignment=PP_ALIGN.CENTER)

    # Subtitle
    tf2 = add_textbox(slide, Inches(2.5), Inches(3.6), Inches(8.333), Inches(0.6))
    set_first_paragraph(tf2, "Using Deep Learning and Computer Vision",
                        font_size=20, color=TEXT_MUTED, font_name=FONT_BODY,
                        alignment=PP_ALIGN.CENTER)

    # Bottom accent line
    add_accent_bar(slide, Inches(5), Inches(4.4), Inches(3.333), Inches(0.03), ACCENT_BLUE)

    # Authors
    tf3 = add_textbox(slide, Inches(2), Inches(4.8), Inches(9.333), Inches(1.0))
    set_first_paragraph(tf3, "Vinoparkavi D  •  Mythily V  •  Dharun Raj R  •  Deepika Y  •  Nithwin V M",
                        font_size=14, color=TEXT_MUTED, font_name=FONT_BODY,
                        alignment=PP_ALIGN.CENTER)
    add_styled_text(tf3, "Department of Computer Science and Engineering",
                    font_size=13, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)
    add_styled_text(tf3, "Nandha Engineering College, Erode, India",
                    font_size=13, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

    # Conference badge
    tf4 = add_textbox(slide, Inches(4.5), Inches(6.3), Inches(4.333), Inches(0.5))
    set_first_paragraph(tf4, "IEEE CONFERENCE PRESENTATION",
                        font_size=12, color=ACCENT_CYAN, bold=True,
                        font_name=FONT_BODY, alignment=PP_ALIGN.CENTER)

    add_slide_number(slide, 1)

def build_intro_slide(prs):
    """Slide 2: Background & Introduction."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_fade_transition(slide)
    add_section_title(slide, "Background & Introduction",
                      "Why automated identity verification matters")

    bullets = [
        "▸  Identity verification is a critical security checkpoint across banking,",
        "    education, and government institutions worldwide.",
        "",
        "▸  Manual verification remains the dominant protocol despite its known",
        "    inefficiencies and error-prone nature.",
        "",
        "▸  Rapid advancements in AI, particularly Convolutional Neural Networks",
        "    (CNNs), have demonstrated superhuman pattern recognition capabilities.",
        "",
        "▸  OCR systems have evolved to extract text from challenging sources,",
        "    including damaged documents and micro-printed text.",
    ]
    add_bullet_block(slide, bullets, MARGIN_LEFT, CONTENT_TOP, Inches(7), start_size=16)

    add_image_safe(slide, "pipeline", Inches(8.2), Inches(2.0), width=Inches(4.5))
    add_slide_number(slide, 2)

def build_problem_slide(prs):
    """Slide 3: The Challenge."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_fade_transition(slide)
    add_section_title(slide, "The Challenge: Manual ID Verification",
                      "Critical bottlenecks in current systems")

    # Stats cards
    stats = [
        ("2–5 min", "Average verification\ntime per person", ACCENT_CYAN),
        ("15–20%", "Human error rate\nunder high volume", ACCENT_ORANGE),
        ("90–180s", "Peak-hour delay\nper individual", ACCENT_PURPLE),
        ("Variable", "Inconsistent judgment\nacross staff", ACCENT_BLUE),
    ]

    for i, (value, label, color) in enumerate(stats):
        x = Inches(0.8 + i * 3.1)
        y = Inches(2.0)

        # Card background
        card = slide.shapes.add_shape(1, x, y, Inches(2.7), Inches(2.0))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor.from_string(BG_CARD)
        card.line.color.rgb = RGBColor.from_string(color)
        card.line.width = Pt(1.5)

        # Accent top bar
        add_accent_bar(slide, x, y, Inches(2.7), Inches(0.04), color)

        # Value
        tf = add_textbox(slide, x + Inches(0.2), y + Inches(0.3), Inches(2.3), Inches(0.7))
        set_first_paragraph(tf, value, font_size=28, color=color,
                            bold=True, alignment=PP_ALIGN.CENTER)

        # Label
        tf2 = add_textbox(slide, x + Inches(0.2), y + Inches(1.1), Inches(2.3), Inches(0.8))
        set_first_paragraph(tf2, label, font_size=13, color=TEXT_MUTED,
                            alignment=PP_ALIGN.CENTER)

    # Bottom note
    bullets = [
        "▸  Errors go both ways — genuine IDs get rejected, fake IDs get accepted.",
        "▸  Security risks compound with fatigue, poor lighting, and inconsistent judgment.",
    ]
    add_bullet_block(slide, bullets, MARGIN_LEFT, Inches(4.5), CONTENT_WIDTH, start_size=15, color=TEXT_MUTED)
    add_slide_number(slide, 3)

def build_limitations_slide(prs):
    """Slide 4: Limitations of Commercial Approaches."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_fade_transition(slide)
    add_section_title(slide, "Limitations of Existing Solutions",
                      "Why current commercial systems fall short")

    items = [
        ("☁️  Cloud Dependency", "Most solutions require constant internet connectivity,\nmaking them unreliable in areas with poor network infrastructure."),
        ("💰  High Operational Costs", "Specialized hardware, premium APIs, and per-use pricing\ncreate barriers for smaller organizations with limited budgets."),
        ("📋  No Universal Standard", "ID cards vary dramatically in size, orientation, material,\nand layout across states, organizations, and countries."),
        ("🔆  Environmental Sensitivity", "Performance degrades severely under inconsistent lighting,\nodd camera angles, and background noise."),
    ]

    for i, (title, desc) in enumerate(items):
        y = Inches(1.8 + i * 1.35)
        # Item accent dot
        add_accent_bar(slide, Inches(1.0), y + Inches(0.12), Inches(0.08), Inches(0.08), ACCENT_CYAN)

        tf = add_textbox(slide, Inches(1.3), y, Inches(10.5), Inches(1.2))
        set_first_paragraph(tf, title, font_size=18, color=TEXT_WHITE, bold=True)
        add_styled_text(tf, desc, font_size=14, color=TEXT_MUTED, spacing_after=Pt(4))

    add_slide_number(slide, 4)

def build_solution_slide(prs):
    """Slide 5: Proposed Solution Overview."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_fade_transition(slide)
    add_section_title(slide, "Our Proposed Solution",
                      "A fully offline, multi-stage AI verification framework")

    bullets = [
        "▸  Integrated Multi-Stage Pipeline combining YOLOv5",
        "    detection, template-based OCR, and FaceNet recognition.",
        "",
        "▸  Runs entirely offline on standard computing",
        "    hardware — no cloud, no special equipment needed.",
        "",
        "▸  Processes complete verifications in under 2 seconds",
        "    with accuracy above 94% across all components.",
        "",
        "▸  Handles 12 different ID card formats with automatic",
        "    lighting adaptation and built-in anti-spoofing.",
    ]
    add_bullet_block(slide, bullets, MARGIN_LEFT, CONTENT_TOP, Inches(6.5), start_size=16)

    add_image_safe(slide, "yolo", Inches(7.8), Inches(1.8), width=Inches(4.8))
    add_slide_number(slide, 5)

def build_architecture_slide(prs):
    """Slide 6: System Architecture & Workflow."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_fade_transition(slide)
    add_section_title(slide, "System Architecture & Workflow",
                      "End-to-end processing pipeline — under 2 seconds")

    # Pipeline image centered
    add_image_safe(slide, "pipeline", Inches(0.8), Inches(1.6), width=Inches(5.5))

    # Right side steps
    steps = [
        ("01", "Scene Capture", "Live webcam video stream analyzed frame by frame"),
        ("02", "Card Detection", "YOLOv5 identifies and extracts ID card region at 25-30 FPS"),
        ("03", "Correction", "Perspective correction and image enhancement applied"),
        ("04", "Parallel OCR + Face", "Template-based text extraction runs alongside FaceNet"),
        ("05", "Final Verdict", "Pass/Fail decision with confidence scores displayed via GUI"),
    ]

    for i, (num, title, desc) in enumerate(steps):
        y = Inches(1.7 + i * 1.05)
        x = Inches(7.0)

        # Step number
        tf_num = add_textbox(slide, x, y, Inches(0.5), Inches(0.4))
        set_first_paragraph(tf_num, num, font_size=14, color=ACCENT_CYAN, bold=True)

        # Title + desc
        tf = add_textbox(slide, x + Inches(0.6), y, Inches(5), Inches(0.9))
        set_first_paragraph(tf, title, font_size=16, color=TEXT_WHITE, bold=True)
        add_styled_text(tf, desc, font_size=12, color=TEXT_MUTED, spacing_after=Pt(2))

    add_slide_number(slide, 6)

def build_dataset_slide(prs):
    """Slide 7: Methodology — Dataset & Training."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_fade_transition(slide)
    add_section_title(slide, "Dataset Construction & Training",
                      "Building a robust detection model")

    # Two-column layout with stat cards
    left_items = [
        ("15,000", "Total training images"),
        ("8,000", "Synthetic templates"),
        ("7,000", "Real campus captures"),
    ]

    right_items = [
        ("100", "Training epochs"),
        ("640×640", "Input resolution"),
        ("18 hrs", "Training duration"),
    ]

    for i, (val, label) in enumerate(left_items):
        y = Inches(1.8 + i * 1.5)
        card = slide.shapes.add_shape(1, Inches(0.8), y, Inches(5.5), Inches(1.2))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor.from_string(BG_CARD)
        card.line.color.rgb = RGBColor.from_string(ACCENT_BLUE)
        card.line.width = Pt(0.75)

        tf = add_textbox(slide, Inches(1.0), y + Inches(0.15), Inches(2), Inches(0.5))
        set_first_paragraph(tf, val, font_size=28, color=ACCENT_CYAN, bold=True)

        tf2 = add_textbox(slide, Inches(3.2), y + Inches(0.15), Inches(3), Inches(0.9))
        set_first_paragraph(tf2, label, font_size=16, color=TEXT_MUTED)

    for i, (val, label) in enumerate(right_items):
        y = Inches(1.8 + i * 1.5)
        card = slide.shapes.add_shape(1, Inches(7.0), y, Inches(5.5), Inches(1.2))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor.from_string(BG_CARD)
        card.line.color.rgb = RGBColor.from_string(ACCENT_PURPLE)
        card.line.width = Pt(0.75)

        tf = add_textbox(slide, Inches(7.2), y + Inches(0.15), Inches(2), Inches(0.5))
        set_first_paragraph(tf, val, font_size=28, color=ACCENT_PURPLE, bold=True)

        tf2 = add_textbox(slide, Inches(9.4), y + Inches(0.15), Inches(3), Inches(0.9))
        set_first_paragraph(tf2, label, font_size=16, color=TEXT_MUTED)

    # Augmentation note
    tf_note = add_textbox(slide, MARGIN_LEFT, Inches(6.3), CONTENT_WIDTH, Inches(0.5))
    set_first_paragraph(tf_note, "Data Augmentation: ±15° rotation  •  0.8–1.2× zoom  •  ±30% brightness  •  Mosaic",
                        font_size=13, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

    add_slide_number(slide, 7)

def build_yolo_slide(prs):
    """Slide 8: Methodology — YOLOv5 Detection."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_fade_transition(slide)
    add_section_title(slide, "Card Detection: YOLOv5",
                      "Real-time object detection for ID card localization")

    bullets = [
        "▸  YOLOv5 chosen for optimal balance of real-time speed",
        "    and detection precision — single-pass architecture.",
        "",
        "▸  Trained for 100 epochs with Adam optimizer",
        "    (β₁=0.9, β₂=0.999), learning rate 0.001.",
        "",
        "▸  Final Performance Metrics:",
        "    •  mAP@0.5:  96.8%",
        "    •  Precision: 94.7%",
        "    •  Recall:    97.2%",
        "",
        "▸  Processing speed: 25–30 FPS on standard hardware.",
    ]
    add_bullet_block(slide, bullets, MARGIN_LEFT, CONTENT_TOP, Inches(6.5), start_size=16)

    add_image_safe(slide, "yolo", Inches(7.8), Inches(1.8), width=Inches(4.8))
    add_slide_number(slide, 8)

def build_ocr_slide(prs):
    """Slide 9: Methodology — Template-Based OCR."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_fade_transition(slide)
    add_section_title(slide, "Template-Based OCR Extraction",
                      "Extensible JSON-driven text recognition system")

    bullets = [
        "▸  Innovative JSON configuration system replaces",
        "    hardcoded ROI coordinates — fully extensible.",
        "",
        "▸  New card types added in under 10 minutes",
        "    without any programming knowledge required.",
        "",
        "▸  Tesseract 4.1 with LSTM neural networks for",
        "    substantially improved character-level accuracy.",
        "",
        "▸  Built-in regex validation automatically flags",
        "    OCR results that don't match expected formats.",
    ]
    add_bullet_block(slide, bullets, MARGIN_LEFT, CONTENT_TOP, Inches(6.5), start_size=16)

    add_image_safe(slide, "ocr", Inches(7.8), Inches(1.8), width=Inches(4.8))
    add_slide_number(slide, 9)

def build_face_slide(prs):
    """Slide 10: Methodology — Facial Recognition."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_fade_transition(slide)
    add_section_title(slide, "Facial Recognition: FaceNet",
                      "128-dimensional embedding-based identity matching")

    bullets = [
        "▸  Pre-trained FaceNet on VGGFace2 dataset —",
        "    millions of facial images for robust performance.",
        "",
        "▸  Converts facial images into 128-D mathematical",
        "    embeddings where similar faces cluster together.",
        "",
        "▸  Calibrated similarity threshold at 0.6 — balancing",
        "    security (2.1% FAR) against usability (3.7% FRR).",
        "",
        "▸  98.2% accuracy validated on Labelled Faces",
        "    in the Wild (LFW) benchmark dataset.",
    ]
    add_bullet_block(slide, bullets, MARGIN_LEFT, CONTENT_TOP, Inches(6.5), start_size=16)

    add_image_safe(slide, "face", Inches(7.8), Inches(1.8), width=Inches(4.8))
    add_slide_number(slide, 10)

def build_antispoofing_slide(prs):
    """Slide 11: Security — Anti-Spoofing."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_fade_transition(slide)
    add_section_title(slide, "Two-Tier Anti-Spoofing Mechanisms",
                      "Multi-layered defense against fraudulent reproductions")

    bullets = [
        "▸  Tier 1 — Frequency Domain Analysis:",
        "    Uses Fourier Transform to detect printer grid patterns",
        "    and halftone screening artifacts invisible to the human eye.",
        "",
        "▸  Tier 2 — Passive Liveness Detection:",
        "    Applies Local Binary Patterns (LBP) to differentiate real",
        "    human skin micro-textures from photographic reproductions.",
        "",
        "▸  Combined efficacy: 85.8% detection of fraud attempts",
        "    with only 4.2% false positive rate on genuine cards.",
    ]
    add_bullet_block(slide, bullets, MARGIN_LEFT, CONTENT_TOP, Inches(6.2), start_size=16)

    add_image_safe(slide, "spoof", Inches(7.5), Inches(1.8), width=Inches(5.0))
    add_slide_number(slide, 11)

def build_results_slide(prs):
    """Slide 12: Experimental Results."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_fade_transition(slide)
    add_section_title(slide, "Experimental Results & Evaluation",
                      "Comprehensive testing across multiple operational scenarios")

    # Left side: key metrics
    metrics = [
        ("98.4%", "Indoor detection accuracy", ACCENT_CYAN),
        ("92.0%", "Challenging outdoor accuracy", ACCENT_BLUE),
        ("94.2%", "OCR text extraction precision", ACCENT_PURPLE),
        ("98.2%", "Face recognition (LFW benchmark)", ACCENT_GREEN),
        ("85.8%", "Anti-spoofing detection rate", ACCENT_ORANGE),
    ]

    for i, (val, label, color) in enumerate(metrics):
        y = Inches(1.7 + i * 0.95)
        add_accent_bar(slide, Inches(0.8), y + Inches(0.1), Inches(0.08), Inches(0.08), color)
        tf = add_textbox(slide, Inches(1.1), y, Inches(1.5), Inches(0.4))
        set_first_paragraph(tf, val, font_size=20, color=color, bold=True)
        tf2 = add_textbox(slide, Inches(2.7), y, Inches(4), Inches(0.4))
        set_first_paragraph(tf2, label, font_size=15, color=TEXT_MUTED)

    # Right side: chart image
    add_image_safe(slide, "results", Inches(7.5), Inches(1.8), width=Inches(5.0))
    add_slide_number(slide, 12)

def build_speed_slide(prs):
    """Slide 13: Processing Speed & Comparison."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_fade_transition(slide)
    add_section_title(slide, "Processing Speed & Comparative Analysis",
                      "Performance vs. commercial and baseline solutions")

    # Comparison table header
    headers = ["Metric", "Our System", "Cloud APIs", "OpenCV Baseline"]
    col_x = [Inches(0.8), Inches(3.8), Inches(6.8), Inches(9.8)]
    col_w = [Inches(2.8), Inches(2.8), Inches(2.8), Inches(2.8)]

    # Header row
    y_header = Inches(2.0)
    add_accent_bar(slide, Inches(0.8), y_header, Inches(11.8), Inches(0.55), ACCENT_BLUE)
    for j, h in enumerate(headers):
        tf = add_textbox(slide, col_x[j], y_header + Inches(0.05), col_w[j], Inches(0.45))
        set_first_paragraph(tf, h, font_size=14, color=TEXT_WHITE, bold=True,
                            alignment=PP_ALIGN.CENTER)

    # Data rows
    rows = [
        ["Detection Accuracy", "96.8%", "97.5%", "84.3%"],
        ["OCR Precision", "94.2%", "96.1%", "87.1%"],
        ["Processing Speed", "25–30 FPS", "Network-dependent", "15–20 FPS"],
        ["Internet Required", "No ✓", "Yes ✗", "No ✓"],
        ["Hardware Cost", "Standard PC", "Cloud + Premium", "Standard PC"],
        ["Anti-Spoofing", "85.8%", "Varies", "None"],
    ]

    for i, row in enumerate(rows):
        y = Inches(2.65 + i * 0.6)
        bg_color = BG_CARD if i % 2 == 0 else BG_DARK
        # Row background
        row_bg = slide.shapes.add_shape(1, Inches(0.8), y, Inches(11.8), Inches(0.55))
        row_bg.fill.solid()
        row_bg.fill.fore_color.rgb = RGBColor.from_string(bg_color)
        row_bg.line.fill.background()

        for j, cell in enumerate(row):
            color = ACCENT_CYAN if j == 1 else TEXT_MUTED
            tf = add_textbox(slide, col_x[j], y + Inches(0.05), col_w[j], Inches(0.45))
            set_first_paragraph(tf, cell, font_size=13, color=color,
                                bold=(j == 0), alignment=PP_ALIGN.CENTER)

    add_slide_number(slide, 13)

def build_deployment_slide(prs):
    """Slide 14: Field Deployment Validation."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_fade_transition(slide)
    add_section_title(slide, "Real-World Field Deployment",
                      "Two-week pilot at Nandha Engineering College")

    bullets = [
        "▸  847 student verifications processed during",
        "    the college registration period.",
        "",
        "▸  92.3% automated verification success rate —",
        "    782 students verified without manual intervention.",
        "",
        "▸  Average processing time: 1.9 seconds per student",
        "    (vs. 90–180 seconds for manual inspection).",
        "",
        "▸  Only 1.5% true technical failures requiring rescan.",
        "    Estimated 8.5 hours of staff time saved.",
    ]
    add_bullet_block(slide, bullets, MARGIN_LEFT, CONTENT_TOP, Inches(6.2), start_size=16)

    add_image_safe(slide, "deployment", Inches(7.5), Inches(1.6), width=Inches(5.2))
    add_slide_number(slide, 14)

def build_conclusion_slide(prs):
    """Slide 15: Conclusion & Future Scope."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_fade_transition(slide)

    # Centered accent bar
    add_accent_bar(slide, Inches(5.5), Inches(0.8), Inches(2.333), Inches(0.03), ACCENT_CYAN)

    # Title
    tf = add_textbox(slide, Inches(1), Inches(1.0), Inches(11.333), Inches(0.7))
    set_first_paragraph(tf, "Conclusion & Future Scope", font_size=32,
                        color=TEXT_WHITE, bold=True, font_name=FONT_TITLE,
                        alignment=PP_ALIGN.CENTER)

    add_accent_bar(slide, Inches(5.5), Inches(1.8), Inches(2.333), Inches(0.02), ACCENT_BLUE)

    # Key contributions
    tf2 = add_textbox(slide, Inches(0.8), Inches(2.2), Inches(6), Inches(0.5))
    set_first_paragraph(tf2, "KEY CONTRIBUTIONS", font_size=14, color=ACCENT_CYAN, bold=True)

    contributions = [
        "▸  Demonstrated a highly accurate, fully offline, and affordable",
        "    verification pipeline running on generic hardware.",
        "",
        "▸  Template-based OCR architecture allows rapid system expansion",
        "    without programming — new cards added in under 10 minutes.",
        "",
        "▸  Multi-layer anti-spoofing identifies 85.8% of fraud attempts.",
    ]
    add_bullet_block(slide, contributions, Inches(0.8), Inches(2.7), Inches(6), start_size=15)

    # Future scope
    tf3 = add_textbox(slide, Inches(7.5), Inches(2.2), Inches(5), Inches(0.5))
    set_first_paragraph(tf3, "FUTURE DIRECTIONS", font_size=14, color=ACCENT_PURPLE, bold=True)

    future = [
        "▸  Upgrade to Transformer-based OCR architectures",
        "    for enhanced text recognition accuracy.",
        "",
        "▸  Implement 3D depth analysis for advanced",
        "    anti-spoofing beyond texture analysis.",
        "",
        "▸  Community-contributed JSON templates for",
        "    expanded international card format support.",
    ]
    add_bullet_block(slide, future, Inches(7.5), Inches(2.7), Inches(5), start_size=15)

    # Thank you
    add_accent_bar(slide, Inches(4), Inches(5.8), Inches(5.333), Inches(0.03), ACCENT_CYAN)
    tf4 = add_textbox(slide, Inches(2), Inches(6.0), Inches(9.333), Inches(0.7))
    set_first_paragraph(tf4, "Thank You — Questions Welcome",
                        font_size=24, color=TEXT_WHITE, bold=True,
                        alignment=PP_ALIGN.CENTER)
    add_styled_text(tf4, "dharunraju7@gmail.com  •  vmnithwin@gmail.com",
                    font_size=13, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

    add_slide_number(slide, 15)


# ═══════════════════════════════════════════════════════════════
#  MAIN EXECUTION
# ═══════════════════════════════════════════════════════════════

def main():
    prs = Presentation()

    # Set widescreen 16:9
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Build all 15 slides
    build_title_slide(prs)       # 1
    build_intro_slide(prs)       # 2
    build_problem_slide(prs)     # 3
    build_limitations_slide(prs) # 4
    build_solution_slide(prs)    # 5
    build_architecture_slide(prs)# 6
    build_dataset_slide(prs)     # 7
    build_yolo_slide(prs)        # 8
    build_ocr_slide(prs)         # 9
    build_face_slide(prs)        # 10
    build_antispoofing_slide(prs)# 11
    build_results_slide(prs)     # 12
    build_speed_slide(prs)       # 13
    build_deployment_slide(prs)  # 14
    build_conclusion_slide(prs)  # 15

    output_path = "/home/shadow/Desktop/Smart_ID_Card_Detection/IEEE_Smart_ID_Card_Presentation.pptx"
    prs.save(output_path)
    print(f"✅ Premium presentation generated successfully with 15 slides!")
    print(f"📁 Saved to: {output_path}")

if __name__ == "__main__":
    main()
