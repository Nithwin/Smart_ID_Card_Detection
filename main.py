# main.py
import argparse
from pptx import Presentation
from pptx.util import Inches, Pt
from content_data import SLIDES_CONTENT
from presentation_utils import (
    download_image, set_slide_background_color, 
    apply_fade_transition, format_title_text, 
    add_bullets_to_frame, PALETTE
)

def create_presentation(output_path):
    prs = Presentation()
    
    # Store layouts
    title_layout = prs.slide_layouts[0]
    content_layout = prs.slide_layouts[1]
    
    for i, slide_data in enumerate(SLIDES_CONTENT):
        layout = title_layout if slide_data["type"] == "title" else content_layout
        slide = prs.slides.add_slide(layout)
        
        # Apply standard background and aesthetic XML transition
        set_slide_background_color(slide, PALETTE["bg"])
        apply_fade_transition(slide)
        
        if slide_data["type"] == "title":
            # Title handling
            t_shape = slide.shapes.title
            body_shape = slide.placeholders[1]
            t_shape.top = Inches(2.5)
            format_title_text(t_shape, slide_data["title"])
            
            # Format subtitle text
            body_shape.text = "\n".join(slide_data["content"])
            for p in body_shape.text_frame.paragraphs:
                p.alignment = 2 # Center alignment
                for run in p.runs:
                    run.font.color.rgb = from_hex(PALETTE["text_light"])
        else:
            # Standard content page
            t_shape = slide.shapes.title
            format_title_text(t_shape, slide_data["title"])
            
            # Adding robust body content
            tf = slide.placeholders[1].text_frame
            add_bullets_to_frame(tf, slide_data["content"])
            
            # Attaching images if specified (only dynamically loads if URL provided)
            # Fits neatly to the right
            if "image" in slide_data and slide_data["image"]:
                img_path = download_image(slide_data["image"], f"slide_{i}.jpg")
                if img_path:
                    slide.shapes.add_picture(img_path, Inches(6.5), Inches(2), width=Inches(3))
                    
    prs.save(output_path)
    print(f"Presentation successfully created with {len(SLIDES_CONTENT)} slides at: {output_path}")

def from_hex(hex_str):
    from pptx.dml.color import RGBColor
    return RGBColor.from_string(hex_str)

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Generate Smart ID Card PPT")
    parser.add_argument("--output", default="IEEE_Smart_ID_Card_Presentation.pptx", help="Path of the generated PPTX")
    args = parser.parse_args()
    create_presentation(args.output)
